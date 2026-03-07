"""Editable PPTX export for DB analysis rounds."""

from __future__ import annotations

import math
from datetime import datetime, timezone
from typing import Iterable

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

from models import DbAnalysisRound


def _page_size_inches(aspect_ratio: str = '16:9', base: float = 10.0) -> tuple[float, float]:
    try:
        w, h = (float(x) for x in aspect_ratio.split(':'))
        if not (math.isfinite(w) and math.isfinite(h) and w > 0 and h > 0):
            raise ValueError
    except Exception:
        w, h = 16.0, 9.0

    if w >= h:
        return base, base * h / w
    return base * w / h, base


class DbAnalysisExportService:
    """Build editable PPTX from structured DB analysis rounds."""

    @staticmethod
    def create_editable_pptx(rounds: Iterable[DbAnalysisRound], output_file: str, aspect_ratio: str = '16:9') -> None:
        prs = Presentation()
        width, height = _page_size_inches(aspect_ratio)
        prs.slide_width = Inches(width)
        prs.slide_height = Inches(height)

        core = prs.core_properties
        now = datetime.now(timezone.utc)
        core.author = 'banana-slides'
        core.last_modified_by = 'banana-slides'
        core.created = now
        core.modified = now

        for round_obj in rounds:
            slide = prs.slides.add_slide(prs.slide_layouts[6])

            # Title
            title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(width - 1.0), Inches(0.7))
            tf = title_box.text_frame
            tf.clear()
            p = tf.paragraphs[0]
            p.text = round_obj.page_title
            p.font.bold = True
            p.font.size = Pt(26)
            p.font.color.rgb = RGBColor(30, 30, 30)
            p.alignment = PP_ALIGN.LEFT

            result_data = round_obj.get_query_result()
            columns = list(result_data.get('columns') or [])[:8]
            rows = list(result_data.get('rows') or [])[:20]

            table_top = Inches(1.2)
            table_height = Inches(3.6)

            if columns:
                table_shape = slide.shapes.add_table(len(rows) + 1, len(columns), Inches(0.5), table_top, Inches(width - 1.0), table_height)
                table = table_shape.table

                for idx, col in enumerate(columns):
                    cell = table.cell(0, idx)
                    cell.text = str(col)
                    cell.text_frame.paragraphs[0].font.bold = True
                    cell.text_frame.paragraphs[0].font.size = Pt(12)

                for r_idx, row_data in enumerate(rows, start=1):
                    for c_idx, col in enumerate(columns):
                        value = row_data.get(col)
                        text_value = '' if value is None else str(value)
                        if len(text_value) > 80:
                            text_value = text_value[:77] + '...'
                        table.cell(r_idx, c_idx).text = text_value
                        table.cell(r_idx, c_idx).text_frame.paragraphs[0].font.size = Pt(11)
            else:
                empty_box = slide.shapes.add_textbox(Inches(0.5), table_top, Inches(width - 1.0), Inches(1.2))
                empty_tf = empty_box.text_frame
                empty_tf.text = '该页无可展示数据表格。'
                empty_tf.paragraphs[0].font.size = Pt(14)

            findings_box = slide.shapes.add_textbox(Inches(0.5), Inches(5.1), Inches(width - 1.0), Inches(1.8))
            findings_tf = findings_box.text_frame
            findings_tf.word_wrap = True
            findings_tf.clear()

            header = findings_tf.paragraphs[0]
            header.text = '关键结论'
            header.font.bold = True
            header.font.size = Pt(16)

            findings = (round_obj.key_findings or '').strip() or '- 无明确结论。'
            for line in findings.splitlines():
                line = line.strip()
                if not line:
                    continue
                para = findings_tf.add_paragraph()
                para.text = line
                para.level = 0
                para.font.size = Pt(12)

        prs.save(output_file)
